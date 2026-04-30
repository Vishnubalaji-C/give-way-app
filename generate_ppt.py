import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_slide(prs, title_text, content_list, is_title_slide=False):
    slide_layout = prs.slide_layouts[1] if not is_title_slide else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Bright background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250) # Light human-made bright background

    title = slide.shapes.title
    title.text = title_text
    
    # Format Title
    title.text_frame.paragraphs[0].font.name = 'Times New Roman'
    title.text_frame.paragraphs[0].font.size = Pt(14)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    if not is_title_slide and len(slide.shapes.placeholders) > 1:
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear() # Clear default
        for item in content_list:
            p = tf.add_paragraph()
            # Split by colon to bold the first part
            parts = item.split(":", 1)
            
            run = p.add_run()
            run.text = parts[0] + (":" if len(parts) > 1 else "")
            run.font.name = 'Times New Roman'
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = RGBColor(51, 51, 51)
            
            if len(parts) > 1:
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.name = 'Times New Roman'
                run2.font.size = Pt(12)
                run2.font.color.rgb = RGBColor(51, 51, 51)
            
            p.space_after = Pt(10)
    elif is_title_slide and len(slide.shapes.placeholders) > 1:
        subtitle = slide.shapes.placeholders[1]
        tf = subtitle.text_frame
        tf.clear()
        for item in content_list:
            p = tf.add_paragraph()
            parts = item.split(":", 1)
            
            run = p.add_run()
            run.text = parts[0] + (":" if len(parts) > 1 else "")
            run.font.name = 'Times New Roman'
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = RGBColor(51, 51, 51)
            
            if len(parts) > 1:
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.name = 'Times New Roman'
                run2.font.size = Pt(12)
                run2.font.color.rgb = RGBColor(51, 51, 51)

prs = Presentation()

slides_data = [
    {
        "title": "GiveWay: Adaptive Traffic Equity System (ATES) utilizing PCE-Weighted Dynamic Phasing",
        "content": [
            "Presenters/Authors: Hari Haran T, Siddharth S, Vishnu Balaji C",
            "Institution: NPR College of Engineering and Technology",
            "Project Overview: A real-time, hardware-software integrated prototype for dynamic traffic signal control based on Passenger Car Equivalent (PCE) metrics.",
            "Aim of the Presentation: To present the architectural design, hardware integration, and operational logic of the GiveWay prototype system."
        ],
        "is_title": True
    },
    {
        "title": "Introduction to the Problem Context",
        "content": [
            "Urban Traffic Congestion: Rapid urbanization has led to increased vehicular density, rendering fixed-time traffic signal controllers inefficient for fluctuating traffic patterns.",
            "Limitations of Static Timers: Pre-programmed traffic lights allocate equal green time regardless of the actual lane density, leading to empty lanes receiving green signals while congested lanes wait.",
            "The Need for Adaptability: Modern intersections require systems that can read current traffic states and dynamically allocate right-of-way to optimize throughput.",
            "Focus on Equity: Traffic management must account for the type of vehicles (e.g., a bus vs. a motorcycle) rather than just the raw count, to ensure true throughput equity."
        ]
    },
    {
        "title": "Proposed Solution: The GiveWay System",
        "content": [
            "Core Concept: GiveWay is an Adaptive Traffic Equity System (ATES) designed to replace static timers with dynamic, real-time phase allocation.",
            "PCE-Centric Approach: Utilizes Passenger Car Equivalent (PCE) weighting to calculate the actual spatial and temporal footprint of vehicles at an intersection, rather than simple vehicle counting.",
            "Hybrid Architecture: Combines edge-level hardware control with cloud-ready backend logic to form a responsive, scalable traffic management node.",
            "Prototype Scope: The current iteration focuses on proving the viability of the logic and hardware integration in a controlled, single-junction prototype environment (e.g., Dindigul junction model)."
        ]
    },
    {
        "title": "Key System Objectives",
        "content": [
            "Dynamic Phase Allocation: To construct a system capable of adjusting green light durations continuously based on live lane density inputs.",
            "Hardware-Software Synchronization: To establish robust, low-latency communication between the traffic control hardware (Arduino/ESP32) and the decision-making backend (Node.js).",
            "Emergency Vehicle Prioritization: To integrate hardware interrupts (RFID) that immediately grant right-of-way to critical emergency services.",
            "Operator Override Capabilities: To provide authorized personnel with the ability to impose manual system states (VIP, Festival, All-Stop) via a centralized dashboard."
        ]
    },
    {
        "title": "System Architecture Overview",
        "content": [
            "Master-Slave Topology: The system is built on a Master-Slave architecture where a central controller dictates phase changes to individual lane nodes.",
            "Edge Hardware Layer: Comprises microcontrollers (Arduino Mega) handling direct electrical actuation of signal LEDs and reading local sensor data.",
            "Decision Engine Layer: A Node.js backend operates as the brain, processing incoming lane data, executing the PCE algorithms, and broadcasting state commands.",
            "Application Interface Layer: A React-based web dashboard and a Flutter/Dart mobile application provide real-time visualization and administrative control over the hardware states."
        ]
    },
    {
        "title": "Hardware Implementation Details",
        "content": [
            "Primary Controller: Arduino Mega acts as the central hardware hub, chosen for its extensive GPIO pin availability necessary for multi-lane LED control.",
            "Visual Indicators: Standard Red, Yellow, and Green LED arrays wired in parallel with custom countdown timer displays for each lane.",
            "Vision Sensors: ESP32-CAM modules positioned at each lane approach to capture visual data of the current traffic state.",
            "Emergency Detection: RFID readers integrated at the approach nodes to detect tagged emergency vehicles and trigger hardware-level interrupts."
        ]
    },
    {
        "title": "Software & Communication Stack",
        "content": [
            "Backend Infrastructure: Built on Node.js and Express, providing a lightweight, non-blocking environment ideal for handling concurrent hardware requests.",
            "Communication Protocols: Utilizes RESTful APIs for state configuration and WebSocket/Serial communication for real-time, low-latency telemetry streaming to the hardware.",
            "Frontend Technologies: React.js is used for the web-based tactical dashboard, ensuring modular component design and rapid UI rendering.",
            "Mobile Interface: Dart/Flutter framework utilized for the mobile application, ensuring cross-platform compatibility for on-the-go system monitoring."
        ]
    },
    {
        "title": "Core Logic: The PCE-Weighted Algorithm",
        "content": [
            "What is PCE? Passenger Car Equivalent (PCE) assigns a standardized spatial value to different vehicle classes (e.g., Two-wheeler = 0.5, Car = 1.0, Bus = 3.0).",
            "Density Calculation: The system aggregates the PCE values of all vehicles detected in a lane to determine the total \"Weight\" or \"Density\" of that lane.",
            "Equity over Count: By using PCE, GiveWay ensures that a lane with two buses receives appropriate priority compared to a lane with three motorcycles, reflecting true spatial demand.",
            "Algorithm Execution: The backend continuously recalculates the PCE density for all active lanes to determine the optimal next phase in the cycle."
        ]
    },
    {
        "title": "Dynamic Phase Allocation Mechanism",
        "content": [
            "Continuous Polling: The master controller continuously polls the status of all four lanes at the junction.",
            "Proportional Time Distribution: Base green time is dynamically extended or reduced proportionally to the calculated PCE density of the active lane.",
            "Phase Skipping: If a lane registers a PCE density of zero (empty lane), the algorithm skips its phase entirely, preventing wasted green time.",
            "Maximum Cap Limits: Hardcoded maximum green time caps are enforced to prevent starvation of lower-density lanes, ensuring eventual right-of-way for all approaches."
        ]
    },
    {
        "title": "Special Operational System Modes",
        "content": [
            "VIP Corridor Mode: An administrative override that locks a specific route in a continuous green state to facilitate the uninterrupted movement of high-priority convoys.",
            "Festival/Congestion Mode: A specialized algorithm state triggered during known high-volume events, adjusting base timings and caps to handle sustained, heavy throughput.",
            "All-Stop / Emergency Halt: A global override that immediately forces all lanes to red, freezing the junction in the event of an accident or critical hazard.",
            "Mode Switching: These modes bypass the standard dynamic algorithm and directly actuate the hardware via the backend override API."
        ]
    },
    {
        "title": "Emergency & Priority Routing",
        "content": [
            "Hardware Interrupts: RFID tags placed on emergency vehicles (ambulances, fire engines) trigger an immediate hardware interrupt when read by the lane sensors.",
            "Preemption Sequence: Upon detection, the system initiates a safe transition sequence (Yellow to Red) for currently active lanes.",
            "Targeted Green Phase: The specific lane containing the emergency vehicle is granted an immediate, locked green phase until the vehicle clears the intersection.",
            "State Restoration: Once the RFID tag is no longer detected in the clearance zone, the system seamlessly resumes the standard PCE-weighted dynamic cycle."
        ]
    },
    {
        "title": "Real-time Monitoring and Interfaces",
        "content": [
            "Tactical Web Dashboard: A React-based interface providing administrators with a live, visual replica of the physical intersection's state.",
            "Live Telemetry: Displays current lane statuses, active countdown timers, and the operational mode currently engaged by the backend.",
            "Incident Tracking: Features a \"Ghost-Lane\" monitor and incident logging system to track anomalies or manual overrides.",
            "Mobile Accessibility: The Flutter app mirrors the critical monitoring functions of the web dashboard, providing field operators with live status updates."
        ]
    },
    {
        "title": "System Advantages and Robustness",
        "content": [
            "Edge-Level Reliability: By utilizing Arduino for direct actuation, the physical light switching remains stable even if backend cloud connectivity experiences momentary latency.",
            "Scalable Architecture: The Node.js and REST/Socket stack allows the single-junction logic to be scaled to a multi-junction network without changing the core edge hardware.",
            "Modular Design: The separation of hardware actuation (Arduino), decision logic (Node), and presentation (React) ensures easy maintenance and future feature additions.",
            "Real-world Applicability: The system is designed using readily available, cost-effective components, making it a highly viable prototype for developing nations' traffic infrastructure."
        ]
    },
    {
        "title": "Challenges and Implementation Constraints",
        "content": [
            "Hardware Limitations: Prototype hardware (Arduino/ESP32) possesses inherent processing and memory constraints compared to industrial-grade PLC controllers.",
            "Environmental Factors: Optical sensors (ESP32-CAMs) are susceptible to environmental noise (glare, heavy rain, low light), requiring robust physical enclosures in a real-world scenario.",
            "Algorithm Boundary Cases: Fine-tuning the minimum and maximum green time caps requires extensive field observation to prevent edge-case starvation scenarios.",
            "Future Scope (AI Integration): While the current system relies on direct PCE calculation logic, future iterations could integrate trained machine learning models for predictive traffic routing."
        ]
    },
    {
        "title": "Conclusion",
        "content": [
            "Project Summary: The GiveWay ATES successfully demonstrates a working prototype of a dynamically adaptive traffic control system.",
            "Achieved Goals: We have successfully integrated edge hardware, a centralized decision engine, and real-time monitoring interfaces into a cohesive ecosystem.",
            "Proof of Concept: The prototype validates the concept that PCE-weighted algorithms can be effectively translated into physical hardware actuation.",
            "Final Statement: GiveWay provides a robust, scalable architectural foundation for modernizing urban traffic infrastructure to be more equitable, responsive, and efficient."
        ]
    }
]

for idx, data in enumerate(slides_data):
    create_slide(prs, data["title"], data["content"], is_title_slide=(idx==0))

output_path = "GiveWay_ATES_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
